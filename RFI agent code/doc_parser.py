"""
Word (.docx) and PowerPoint (.pptx) RFI parser.
Extracts Q&A pairs from documents and returns them as RFIQuestion objects
(same structure as excel_parser.py for pipeline compatibility).
"""

from __future__ import annotations

import os
import re
from typing import Optional

from docx import Document
from docx.table import Table
from pptx import Presentation
from pptx.util import Inches

from excel_parser import RFIQuestion, _looks_like_question, _is_section_header


# ─── WORD (.docx) PARSING ────────────────────────────────────────────────────

# Patterns for detecting numbered questions
NUMBERED_Q_PATTERN = re.compile(
    r"^(\d+[\.\)]\s*\d*\.?\s*|[a-zA-Z][\.\)]\s*|Q\d+[\.\):\s])",
    re.IGNORECASE,
)

# Patterns for section headings in Word docs
SECTION_HEADING_PATTERN = re.compile(
    r"^(section\s+\d|part\s+\d|\d+\.\s+[A-Z]|[A-Z]\.\s+[A-Z]|chapter\s+\d)",
    re.IGNORECASE,
)


def _extract_question_number_from_text(text: str) -> tuple[str, str]:
    """
    Extract question number prefix from text.
    Returns (number, remaining_text).
    """
    match = NUMBERED_Q_PATTERN.match(text)
    if match:
        num = match.group(1).strip().rstrip(".):").strip()
        remaining = text[match.end():].strip()
        return num, remaining
    return "", text


def _infer_category_from_heading(heading: str) -> str:
    """Map document heading text to a category hint."""
    lower = heading.lower()
    mappings = {
        "company": "Company Information",
        "overview": "Company Information",
        "general": "Company Information",
        "compliance": "Compliance",
        "legal": "Compliance",
        "ethics": "Compliance",
        "data": "Data & Information Security",
        "security": "Data & Information Security",
        "privacy": "Data & Information Security",
        "cyber": "Data & Information Security",
        "environment": "ESG",
        "sustainability": "ESG",
        "esg": "ESG",
        "social": "ESG",
        "people": "People Information",
        "staff": "People Information",
        "team": "People Information",
        "resource": "People Information",
        "supplier": "Suppliers & Freelancers",
        "subcontract": "Suppliers & Freelancers",
        "vendor": "Suppliers & Freelancers",
        "technolog": "Technology & AI",
        "digital": "Technology & AI",
        "ai": "Technology & AI",
        "commercial": "Commercial Information",
        "financial": "Commercial Information",
        "pricing": "Commercial Information",
        "capabilit": "Commercial Information",
    }
    for keyword, category in mappings.items():
        if keyword in lower:
            return category
    return ""


def _parse_docx_tables(doc: Document) -> list[RFIQuestion]:
    """
    Extract questions from tables in a Word document.
    Many RFIs embed questions in tables with columns like:
    [#, Question, Answer] or [Requirement, Response, Evidence]
    """
    questions = []

    for table_idx, table in enumerate(doc.tables):
        if len(table.rows) < 2:
            continue

        # Detect header row
        header_cells = [cell.text.strip() for cell in table.rows[0].cells]

        # Find question and answer columns
        q_col = None
        a_col = None
        n_col = None

        from excel_parser import QUESTION_HEADERS, ANSWER_HEADERS, NUMBER_HEADERS

        for i, header in enumerate(header_cells):
            if not header:
                continue
            if QUESTION_HEADERS.search(header) and q_col is None:
                q_col = i
            elif ANSWER_HEADERS.search(header) and a_col is None:
                a_col = i
            elif NUMBER_HEADERS.search(header) and n_col is None:
                n_col = i

        # Fallback: if no headers detected, assume col 0 or 1 is questions
        if q_col is None:
            if len(header_cells) >= 3:
                q_col = 1  # [#, Question, Answer]
                if n_col is None:
                    n_col = 0
            else:
                q_col = 0  # [Question, Answer]

        if a_col is None and len(header_cells) > q_col + 1:
            a_col = q_col + 1

        # Extract questions from data rows
        for row_idx, row in enumerate(table.rows[1:], start=2):
            cells = [cell.text.strip() for cell in row.cells]
            if q_col >= len(cells):
                continue

            q_text = cells[q_col]
            a_text = cells[a_col] if a_col is not None and a_col < len(cells) else ""
            q_num = cells[n_col] if n_col is not None and n_col < len(cells) else ""

            if not q_text or not _looks_like_question(q_text):
                continue
            if _is_section_header(q_text):
                continue

            questions.append(RFIQuestion(
                sheet_name=f"Table {table_idx + 1}",
                row=row_idx,
                question_col="A",
                answer_col="B" if a_col is not None else None,
                question_number=q_num,
                question_text=q_text,
                existing_answer=a_text,
                category_hint="",
            ))

    return questions


def _parse_docx_paragraphs(doc: Document) -> list[RFIQuestion]:
    """
    Extract questions from paragraphs (numbered lists, headings, body text).
    Handles formats like:
    - "1. What is your company name?"
    - "Q1: Describe your services"
    - Bullet/numbered lists under section headings
    """
    questions = []
    current_section = ""
    current_category = ""
    row_counter = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Check if this is a heading / section header
        is_heading = para.style.name.startswith("Heading") if para.style else False

        if is_heading or _is_section_header(text) or SECTION_HEADING_PATTERN.match(text):
            current_section = text
            current_category = _infer_category_from_heading(text)
            continue

        # Try to extract question number
        q_num, q_text = _extract_question_number_from_text(text)

        # If we extracted a number, the remaining text is more likely a question
        if q_num and _looks_like_question(q_text):
            row_counter += 1
            questions.append(RFIQuestion(
                sheet_name=current_section or "Document",
                row=row_counter,
                question_col="A",
                answer_col=None,
                question_number=q_num,
                question_text=q_text,
                existing_answer="",
                category_hint=current_category,
            ))
        elif not q_num and _looks_like_question(text):
            # No number prefix but still looks like a question
            row_counter += 1
            questions.append(RFIQuestion(
                sheet_name=current_section or "Document",
                row=row_counter,
                question_col="A",
                answer_col=None,
                question_number="",
                question_text=text,
                existing_answer="",
                category_hint=current_category,
            ))

    return questions


def parse_docx(filepath: str) -> list[RFIQuestion]:
    """
    Parse a Word document RFI and extract all Q&A pairs.
    Tries tables first (most structured), then falls back to paragraph parsing.
    """
    doc = Document(filepath)

    # Try tables first — most RFI Word docs use tables
    table_questions = _parse_docx_tables(doc)
    if table_questions:
        # If tables yielded results, also get paragraph questions for sections
        # not covered by tables
        para_questions = _parse_docx_paragraphs(doc)
        # Deduplicate: if a paragraph question matches a table question, skip it
        table_texts = {q.question_text.lower().strip() for q in table_questions}
        unique_para = [
            q for q in para_questions
            if q.question_text.lower().strip() not in table_texts
        ]
        return table_questions + unique_para

    # No tables — parse paragraphs only
    return _parse_docx_paragraphs(doc)


# ─── POWERPOINT (.pptx) PARSING ──────────────────────────────────────────────


def _parse_pptx_tables(prs: Presentation) -> list[RFIQuestion]:
    """Extract questions from tables embedded in slides."""
    questions = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        # Get slide title for category hint
        slide_title = ""
        if slide.shapes.title:
            slide_title = slide.shapes.title.text.strip()

        category = _infer_category_from_heading(slide_title) if slide_title else ""

        for shape in slide.shapes:
            if not shape.has_table:
                continue

            table = shape.table
            if len(table.rows) < 2:
                continue

            # Detect header row
            header_cells = [table.cell(0, col).text.strip() for col in range(len(table.columns))]

            from excel_parser import QUESTION_HEADERS, ANSWER_HEADERS, NUMBER_HEADERS

            q_col = None
            a_col = None
            n_col = None

            for i, header in enumerate(header_cells):
                if not header:
                    continue
                if QUESTION_HEADERS.search(header) and q_col is None:
                    q_col = i
                elif ANSWER_HEADERS.search(header) and a_col is None:
                    a_col = i
                elif NUMBER_HEADERS.search(header) and n_col is None:
                    n_col = i

            if q_col is None:
                q_col = 1 if len(header_cells) >= 3 else 0
            if a_col is None and len(header_cells) > q_col + 1:
                a_col = q_col + 1

            # Extract questions
            for row_idx in range(1, len(table.rows)):
                q_text = table.cell(row_idx, q_col).text.strip() if q_col < len(table.columns) else ""
                a_text = table.cell(row_idx, a_col).text.strip() if a_col is not None and a_col < len(table.columns) else ""
                q_num = table.cell(row_idx, n_col).text.strip() if n_col is not None and n_col < len(table.columns) else ""

                if not q_text or not _looks_like_question(q_text):
                    continue
                if _is_section_header(q_text):
                    continue

                questions.append(RFIQuestion(
                    sheet_name=slide_title or f"Slide {slide_idx}",
                    row=row_idx + 1,
                    question_col="A",
                    answer_col="B" if a_col is not None else None,
                    question_number=q_num,
                    question_text=q_text,
                    existing_answer=a_text,
                    category_hint=category,
                ))

    return questions


def _parse_pptx_text_frames(prs: Presentation) -> list[RFIQuestion]:
    """Extract questions from slide text frames (bullet points, body text)."""
    questions = []
    row_counter = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_title = ""
        if slide.shapes.title:
            slide_title = slide.shapes.title.text.strip()

        category = _infer_category_from_heading(slide_title) if slide_title else ""

        for shape in slide.shapes:
            if shape.has_table:
                continue  # Tables handled separately
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # Skip the title text itself (already captured)
                if text == slide_title:
                    continue

                # Skip section headers
                if _is_section_header(text):
                    continue

                # Extract question number if present
                q_num, q_text = _extract_question_number_from_text(text)

                if q_num and _looks_like_question(q_text):
                    row_counter += 1
                    questions.append(RFIQuestion(
                        sheet_name=slide_title or f"Slide {slide_idx}",
                        row=row_counter,
                        question_col="A",
                        answer_col=None,
                        question_number=q_num,
                        question_text=q_text,
                        existing_answer="",
                        category_hint=category,
                    ))
                elif not q_num and _looks_like_question(text):
                    row_counter += 1
                    questions.append(RFIQuestion(
                        sheet_name=slide_title or f"Slide {slide_idx}",
                        row=row_counter,
                        question_col="A",
                        answer_col=None,
                        question_number="",
                        question_text=text,
                        existing_answer="",
                        category_hint=category,
                    ))

    return questions


def parse_pptx(filepath: str) -> list[RFIQuestion]:
    """
    Parse a PowerPoint RFI and extract all Q&A pairs.
    Tries tables first, then text frames.
    """
    prs = Presentation(filepath)

    # Tables first (most structured)
    table_questions = _parse_pptx_tables(prs)
    if table_questions:
        text_questions = _parse_pptx_text_frames(prs)
        table_texts = {q.question_text.lower().strip() for q in table_questions}
        unique_text = [
            q for q in text_questions
            if q.question_text.lower().strip() not in table_texts
        ]
        return table_questions + unique_text

    return _parse_pptx_text_frames(prs)


# ─── UNIFIED ENTRY POINT ─────────────────────────────────────────────────────


def parse_document(filepath: str) -> list[RFIQuestion]:
    """
    Unified parser entry point. Detects format by extension and delegates.
    Supports: .docx, .pptx
    For .xlsx/.xlsm, use excel_parser.parse_rfi() directly.
    """
    ext = os.path.splitext(filepath)[1].lower()

    if ext == ".docx":
        return parse_docx(filepath)
    elif ext == ".pptx":
        return parse_pptx(filepath)
    else:
        raise ValueError(f"Unsupported document format: {ext}. Use .docx or .pptx")


def extract_client_from_filename(filename: str) -> str:
    """Try to extract client name from document filename (same logic as excel_parser)."""
    from excel_parser import extract_client_from_filename as _extract
    return _extract(filename)
