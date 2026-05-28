#!/usr/bin/env python3
"""Extract content from all RFI library files."""

import os
import openpyxl
from docx import Document
from pptx import Presentation

RFI_DIR = os.path.join(os.path.dirname(__file__), "RFI library")

def extract_xlsx(filepath):
    """Extract all text from an xlsx/xlsm file."""
    wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
    result = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        result.append(f"\n=== SHEET: {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            line = " | ".join(cells)
            if line.replace("|", "").strip():
                result.append(line)
    wb.close()
    return "\n".join(result)

def extract_docx(filepath):
    """Extract all text from a docx file."""
    doc = Document(filepath)
    result = []
    for para in doc.paragraphs:
        if para.text.strip():
            result.append(para.text)
    for table in doc.tables:
        result.append("\n--- TABLE ---")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            result.append(" | ".join(cells))
    return "\n".join(result)

def extract_pptx(filepath):
    """Extract all text from a pptx file."""
    prs = Presentation(filepath)
    result = []
    for i, slide in enumerate(prs.slides, 1):
        result.append(f"\n=== SLIDE {i} ===")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                result.append(shape.text)
            if shape.has_table:
                result.append("--- TABLE ---")
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    result.append(" | ".join(cells))
    return "\n".join(result)

def main():
    files = sorted(os.listdir(RFI_DIR))
    for fname in files:
        filepath = os.path.join(RFI_DIR, fname)
        print(f"\n{'='*80}")
        print(f"FILE: {fname}")
        print(f"{'='*80}")
        try:
            if fname.endswith(('.xlsx', '.xlsm')):
                content = extract_xlsx(filepath)
            elif fname.endswith('.docx'):
                content = extract_docx(filepath)
            elif fname.endswith('.pptx'):
                content = extract_pptx(filepath)
            else:
                content = f"[Unsupported format: {fname}]"
            print(content)
        except Exception as e:
            print(f"ERROR: {e}")

if __name__ == "__main__":
    main()
