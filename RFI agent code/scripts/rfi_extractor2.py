#!/usr/bin/env python3
"""Extract content from all RFI library files into separate output files."""

import os
import sys
import openpyxl
from docx import Document
from pptx import Presentation

RFI_DIR = os.path.join(os.path.dirname(__file__), "RFI library")
OUT_DIR = os.path.join(os.path.dirname(__file__), "rfi_extracted")
os.makedirs(OUT_DIR, exist_ok=True)

def extract_xlsx(filepath):
    wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
    result = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        result.append(f"\n=== SHEET: {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            line = " | ".join(cells)
            if line.replace("|", "").strip():
                result.append(line)
    wb.close()
    return "\n".join(result)

def extract_docx(filepath):
    doc = Document(filepath)
    result = []
    for para in doc.paragraphs:
        if para.text.strip():
            result.append(para.text)
    for table in doc.tables:
        result.append("\n--- TABLE ---")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            result.append(" | ".join(cells))
    return "\n".join(result)

def extract_pptx(filepath):
    prs = Presentation(filepath)
    result = []
    for i, slide in enumerate(prs.slides, 1):
        result.append(f"\n=== SLIDE {i} ===")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                result.append(shape.text)
            if shape.has_table:
                result.append("--- TABLE ---")
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    result.append(" | ".join(cells))
    return "\n".join(result)

def main():
    files = sorted(os.listdir(RFI_DIR))
    print(f"Found {len(files)} files")
    for i, fname in enumerate(files, 1):
        filepath = os.path.join(RFI_DIR, fname)
        out_name = f"{i:02d}_{os.path.splitext(fname)[0]}.txt"
        out_path = os.path.join(OUT_DIR, out_name)
        print(f"[{i}/{len(files)}] {fname}")
        try:
            if fname.endswith(('.xlsx', '.xlsm')):
                content = extract_xlsx(filepath)
            elif fname.endswith('.docx'):
                content = extract_docx(filepath)
            elif fname.endswith('.pptx'):
                content = extract_pptx(filepath)
            else:
                content = f"[Unsupported format: {fname}]"
            with open(out_path, 'w', encoding='utf-8') as f:
                f.write(content)
            print(f"  -> {out_name} ({len(content)} chars)")
        except Exception as e:
            print(f"  ERROR: {e}")

if __name__ == "__main__":
    main()
